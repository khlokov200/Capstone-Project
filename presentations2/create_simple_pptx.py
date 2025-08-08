"""
Create a very simple PowerPoint presentation that should work with Keynote
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

# Config
TITLE = "The MeteoMetrics Weather Station"
SUBTITLE = "A professional-grade weather dashboard by Tobi Odika"
SCREEN_DIR = Path("../screens")  # Path relative to the presentations2 folder
OUTPUT = "MeteoMetrics_Simple.pptx"

def create_simple_pptx():
    """Create a very simple PowerPoint that should be compatible with Keynote"""
    prs = Presentation()
    
    # Title slide - use a very simple approach
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if hasattr(slide.shapes, "title") and slide.shapes.title is not None:
        slide.shapes.title.text = TITLE
    
    # Add subtitle manually as a text box
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = SUBTITLE
    
    # Executive Summary - add as a basic title slide with text box
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if hasattr(slide.shapes, "title") and slide.shapes.title is not None:
        slide.shapes.title.text = "Executive Summary"
    
    # Add content as a text box
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "• 15+ functional tabs with advanced weather tools"
    p = tf.add_paragraph()
    p.text = "• Real-time data, analytics, and machine learning"
    p = tf.add_paragraph()
    p.text = "• 6+ chart types, modern UI, and zero critical bugs"
    
    # Add image slides
    images = sorted(SCREEN_DIR.glob("*.png"))
    captions = [
        "Quick Actions Dashboard",
        "Current Weather View",
        "Forecast Tab",
        "Live Weather Radar",
        "Analytics & Trends",
        "City Comparison", 
        "Health & Wellness",
        "Activity Suggestions",
        "Weather Poetry",
        "Weather History",
        "Settings",
        "Extra Feature"
    ]
    
    for img_path, caption in zip(images, captions):
        # Use completely blank slide for maximum compatibility
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add title as text box
        left = Inches(1)
        top = Inches(0.5)
        width = Inches(8)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = caption
        
        # Add image
        try:
            left = Inches(1)
            top = Inches(1.5)
            slide.shapes.add_picture(str(img_path), left, top, height=Inches(5))
        except Exception as e:
            print(f"Error adding image {img_path}: {e}")
    
    # Save presentation
    prs.save(OUTPUT)
    print(f"Created a simple PowerPoint file: {OUTPUT}")

if __name__ == "__main__":
    create_simple_pptx()
