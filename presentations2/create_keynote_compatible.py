"""
Create a presentation in PowerPoint XML format, which is more compatible with Keynote
"""
import os
from pathlib import Path
import xml.etree.ElementTree as ET
import zipfile
import io
import base64
import shutil
import tempfile

# Config
TITLE = "The MeteoMetrics Weather Station"
SUBTITLE = "A professional-grade weather dashboard by Tobi Odika"
SCREEN_DIR = Path("../screens")  # Path relative to the presentations2 folder
OUTPUT = "MeteoMetrics_Keynote_Compatible.pptx"

def create_compatible_pptx():
    """Create a basic PowerPoint file that should be more compatible with Keynote"""
    try:
        # Check if we can run the more compatible PowerPoint generator
        import pip
        pip.main(['install', 'python-pptx==0.6.18'])  # Use an older version for better compatibility
        
        # Now try again with the older version of python-pptx
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
        slide.shapes.title.text = TITLE
        if hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
            slide.placeholders[1].text = SUBTITLE
        
        # Executive Summary slide - use simpler layout with just a title
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
        slide.shapes.title.text = "Executive Summary"
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = "• 15+ functional tabs with advanced weather tools"
            p = tf.add_paragraph()
            p.text = "• Real-time data, analytics, and machine learning"
            p = tf.add_paragraph()
            p.text = "• 6+ chart types, modern UI, and zero critical bugs"
        
        # Key Features slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Key Features"
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = "• Real-time weather & multi-city comparison"
            p = tf.add_paragraph()
            p.text = "• 5-day forecasting & historical trends"
            p = tf.add_paragraph()
            p.text = "• Live radar animations & severe alerts"
            p = tf.add_paragraph()
            p.text = "• Health & wellness monitoring"
            p = tf.add_paragraph()
            p.text = "• AI-based activity & travel suggestions"
        
        # Technical Architecture slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Technical Architecture"
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = "• MVC pattern with modular components"
            p = tf.add_paragraph()
            p.text = "• Tkinter frontend, matplotlib charts"
            p = tf.add_paragraph()
            p.text = "• OpenWeatherMap API integration"
            p = tf.add_paragraph()
            p.text = "• Machine learning for predictions"
        
        # UI Highlights slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "User Interface Highlights"
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = "• Modern split-panel layout"
            p = tf.add_paragraph()
            p.text = "• Professional styling & color schemes"
            p = tf.add_paragraph()
            p.text = "• Interactive charts with multiple types"
            p = tf.add_paragraph()
            p.text = "• Quick Actions dashboard"
        
        # Add image slides - very simple to maximize compatibility
        images = sorted(SCREEN_DIR.glob("*.png"))
        captions = [
            "Quick Actions Dashboard",
            "Current Weather View",
            "Forecast Tab",
            "Live Weather Radar",
            "Analytics & Trends",
            "City Comparison",
            "Health & Wellness",
            "Activity Suggestions",
            "Weather Poetry",
            "Weather History",
            "Settings",
            "Extra Feature"
        ]
        
        for img_path, caption in zip(images, captions):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            slide.shapes.title.text = caption
            try:
                slide.shapes.add_picture(str(img_path), Inches(1), Inches(1.5), height=Inches(5))
            except Exception as e:
                print(f"Error adding image {img_path}: {e}")
                # Add a text placeholder instead
                txbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(5))
                txbox.text_frame.text = f"[Image: {img_path.name}]"
        
        # Testing slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Testing & Quality Assurance"
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = "• All UI components tested"
            p = tf.add_paragraph()
            p.text = "• Chart generation validated"
            p = tf.add_paragraph()
            p.text = "• API integration stable"
            p = tf.add_paragraph()
            p.text = "• No critical bugs"
        
        # Future Enhancements slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Future Enhancements"
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = "• Mobile responsive design"
            p = tf.add_paragraph()
            p.text = "• SQL database integration"
            p = tf.add_paragraph()
            p.text = "• Enhanced ML models"
            p = tf.add_paragraph()
            p.text = "• IoT weather station integration"
        
        # Save presentation
        prs.save(OUTPUT)
        print(f"Created a Keynote-compatible PowerPoint file: {OUTPUT}")
    except Exception as e:
        print(f"Error creating compatible PowerPoint: {e}")
        print("You can still use the HTML presentation which should work in any web browser.")

if __name__ == "__main__":
    create_compatible_pptx()
